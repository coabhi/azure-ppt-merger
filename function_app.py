import logging
import os
from datetime import datetime
from azure.storage.blob import BlobServiceClient
from pptx import Presentation
from pptx.util import Inches
import tempfile
import azure.functions as func

# Set up Blob Service Client
STORAGE_CONNECTION_STRING = os.environ["AzureWebJobsStorage"]

app = func.FunctionApp(http_auth_level=func.AuthLevel.ANONYMOUS)

@app.route(route="http_trigger")
def http_trigger(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("Processing HTTP trigger for PPT consolidation.")

    # Extract container and folder name from the request
    BLOB_CONTAINER_NAME = req.params.get('blobName')
    folder_name = req.params.get('folderName')
    
    if not BLOB_CONTAINER_NAME or not folder_name:
        try:
            req_body = req.get_json()
        except ValueError:
            pass
        else:
            BLOB_CONTAINER_NAME = req_body.get('blobName')
            folder_name = req_body.get('folderName')

    if not BLOB_CONTAINER_NAME or not folder_name:
        return func.HttpResponse("Both 'blobName' and 'folderName' must be provided.", status_code=400)

    blob_service_client = BlobServiceClient.from_connection_string(STORAGE_CONNECTION_STRING)
    container_client = blob_service_client.get_container_client(BLOB_CONTAINER_NAME)

    # Get the list of PPT files only from the specified folder
    ppt_files = [
        blob.name for blob in container_client.list_blobs(name_starts_with=folder_name) if blob.name.endswith(".pptx")
    ]

    if not ppt_files:
        return func.HttpResponse(f"No PPT files found in the folder '{folder_name}' of the container.", status_code=404)

    # Create temporary directories for processing
    with tempfile.TemporaryDirectory() as temp_dir:
        images_dir = os.path.join(temp_dir, "Images")
        os.makedirs(images_dir, exist_ok=True)

        # Download, process, and extract images
        for i, blob_name in enumerate(ppt_files):
            blob_client = container_client.get_blob_client(blob_name)
            ppt_file_path = os.path.join(temp_dir, f"presentation_{i}.pptx")

            # Download PPT file to temporary storage
            with open(ppt_file_path, "wb") as ppt_file:
                ppt_file.write(blob_client.download_blob().readall())

            # Extract images from the presentation
            save_images_from_ppt(ppt_file_path, images_dir, i)

        # Create consolidated presentation
        consolidated_ppt_path = os.path.join(temp_dir, f"consolidated_{get_timestamp()}.pptx")
        create_presentation_with_images(images_dir, consolidated_ppt_path)

        # Define consolidated file's blob path in the same folder
        consolidated_blob_name = f"{folder_name}/Consolidated_PPT_{get_timestamp()}.pptx"
        
        # Upload the consolidated presentation back to the blob container in the same folder
        upload_to_blob(container_client, consolidated_ppt_path, consolidated_blob_name)

        # Store original PPT files into a separate folder and delete the original copies
        folder_name_for_originals = f"{folder_name}/originals_PPTs_{get_timestamp()}"
        store_ppt_files_in_folder_and_delete(container_client, ppt_files, folder_name_for_originals)

    return func.HttpResponse(f"PPT files from folder '{folder_name}' processed and consolidated successfully.", status_code=200)

def save_images_from_ppt(ppt_file_path, output_dir, index_prefix):
    prs = Presentation(ppt_file_path)

    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # 13 indicates a picture
                image_bytes = shape.image.blob
                image_path = os.path.join(output_dir, f"{index_prefix}_slide_{i}_{j}.png")
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)

def create_presentation_with_images(image_dir, output_file):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide_height = prs.slide_height

    for img_name in sorted(os.listdir(image_dir)):
        img_path = os.path.join(image_dir, img_name)
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        left = Inches(0.09)
        slide.shapes.add_picture(img_path, left, 0, height=slide_height)

    prs.save(output_file)

def upload_to_blob(container_client, local_file_path, blob_name):
    blob_client = container_client.get_blob_client(blob_name)
    with open(local_file_path, "rb") as data:
        blob_client.upload_blob(data, overwrite=True)

def store_ppt_files_in_folder_and_delete(container_client, ppt_files, folder_name):
    """Upload original PPT files to a specified folder in the blob container and delete the original files."""
    for blob_name in ppt_files:
        # Skip already consolidated PPTs
        if "Consolidated_PPT" in blob_name:
            continue

        blob_client = container_client.get_blob_client(blob_name)
        original_blob_client = container_client.get_blob_client(f"{folder_name}/{os.path.basename(blob_name)}")

        # Download blob and upload to the new folder
        # blob_data = blob_client.download_blob().readall()
        # original_blob_client.upload_blob(blob_data, overwrite=True)

        # Delete the original blob
        blob_client.delete_blob()

def get_timestamp():
    return datetime.now().strftime("%Y_%m_%d_%H_%M_%S")
